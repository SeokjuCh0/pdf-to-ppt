#!/usr/bin/env python3
"""Machine-readable CLI for the local PDF-to-PPT conversion engine."""

from __future__ import annotations

import argparse
import json
import shutil
import sys
import tempfile
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation

from scripts.extract_pdf_json import extract_pdf_json, find_parser_jar
from scripts.json_to_pptx import convert_json_to_pptx, grouped_by_page, iter_renderable_objects, text_from_object
from scripts.pdf_to_pptx import copy_json_with_external_assets


def emit(payload: dict[str, Any]) -> None:
    print(json.dumps(payload, ensure_ascii=False, indent=2))


def walk(value: Any) -> list[dict[str, Any]]:
    found: list[dict[str, Any]] = []
    if isinstance(value, dict):
        found.append(value)
        for child in value.values():
            found.extend(walk(child))
    elif isinstance(value, list):
        for child in value:
            found.extend(walk(child))
    return found


def summarize_layout(data: dict[str, Any]) -> dict[str, Any]:
    objects = walk(data.get("kids", []))
    types = Counter(str(obj.get("type") or "unknown").lower() for obj in objects)
    text_objects = [
        obj for obj in objects
        if isinstance(obj.get("content"), str) and obj.get("content", "").strip()
    ]
    image_only_charts = types.get("image", 0)
    warnings: list[str] = []
    if image_only_charts:
        warnings.append("Some visual regions were extracted as images, not editable tables/charts.")
    return {
        "pages": int(data.get("number of pages") or 0),
        "objects": len(objects),
        "types": dict(sorted(types.items())),
        "text_objects": len(text_objects),
        "tables": types.get("table", 0),
        "images": types.get("image", 0),
        "warnings": warnings,
    }


def list_components(data: dict[str, Any]) -> list[dict[str, Any]]:
    components: list[dict[str, Any]] = []
    for page in sorted(grouped_by_page(data)):
        page_items = grouped_by_page(data)[page]
        for index, obj in enumerate(iter_renderable_objects(page_items), start=1):
            obj_type = str(obj.get("type") or "unknown").lower()
            bbox = obj.get("bounding box")
            component_id = str(obj.get("id") or f"{page}:{index}")
            text = text_from_object(obj)
            components.append({
                "id": component_id,
                "page": int(obj.get("page number") or page),
                "type": obj_type,
                "editable": obj_type not in {"image"},
                "bounding_box": bbox if isinstance(bbox, list) and len(bbox) == 4 else None,
                "font": obj.get("font"),
                "font_size": obj.get("font size"),
                "text": text,
                "source": obj.get("source"),
                "label": text[:80] if text else obj.get("source") or obj_type,
            })
    return components


def summarize_pptx(pptx_path: Path) -> dict[str, Any]:
    prs = Presentation(pptx_path)
    text_shapes = 0
    tables = 0
    pictures = 0
    shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes += 1
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                text_shapes += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if int(shape.shape_type) == 13:
                pictures += 1
    warnings: list[str] = []
    if pictures and not tables:
        warnings.append("PPTX contains images but no editable tables; chart interiors may be image-only.")
    return {
        "slides": len(prs.slides),
        "shapes": shapes,
        "text_shapes": text_shapes,
        "tables": tables,
        "pictures": pictures,
        "warnings": warnings,
    }


def build_parser_common(parser: argparse.ArgumentParser) -> None:
    parser.add_argument("--jar", help="Repo-built OpenDataLoader parser JAR.")
    parser.add_argument("--java", default="java", help="Java command to run.")
    parser.add_argument("--pages", help='Pages to extract, for example "1,3,5-7".')
    parser.add_argument("--table-method", default="cluster", choices=["default", "cluster"])
    parser.add_argument("--reading-order", default="xycut", choices=["off", "xycut"])
    parser.add_argument("--image-output", default="external", choices=["off", "embedded", "external"])
    parser.add_argument("--work-dir", type=Path, help="Directory for parser outputs.")


def command_inspect(args: argparse.Namespace) -> int:
    jar = find_parser_jar(args.jar)
    pdf = args.pdf.expanduser().resolve()
    with tempfile.TemporaryDirectory(prefix="pdfppt-inspect-") as tmp:
        work_dir = args.work_dir.expanduser().resolve() if args.work_dir else Path(tmp)
        json_path = extract_pdf_json(
            pdf,
            work_dir,
            jar=jar,
            java_cmd=args.java,
            pages=args.pages,
            table_method=args.table_method,
            reading_order=args.reading_order,
            image_output=args.image_output,
        )
        data = json.loads(json_path.read_text(encoding="utf-8"))
        if args.json_output:
            copy_json_with_external_assets(json_path, args.json_output.expanduser().resolve())
        emit({
            "ok": True,
            "command": "inspect",
            "pdf": str(pdf),
            "json": str(args.json_output.expanduser().resolve()) if args.json_output else str(json_path),
            "layout": summarize_layout(data),
            "components": list_components(data),
        })
    return 0


def command_convert(args: argparse.Namespace) -> int:
    jar = find_parser_jar(args.jar)
    pdf = args.pdf.expanduser().resolve()
    pptx = args.pptx.expanduser().resolve()
    if args.work_dir:
        work_dir = args.work_dir.expanduser().resolve()
        work_dir.mkdir(parents=True, exist_ok=True)
        cleanup = None
    else:
        cleanup = tempfile.TemporaryDirectory(prefix="pdfppt-convert-")
        work_dir = Path(cleanup.name)
    try:
        json_path = extract_pdf_json(
            pdf,
            work_dir,
            jar=jar,
            java_cmd=args.java,
            pages=args.pages,
            table_method=args.table_method,
            reading_order=args.reading_order,
            image_output=args.image_output,
        )
        convert_json_to_pptx(
            json_path,
            pptx,
            page_width=args.page_width,
            page_height=args.page_height,
            fallback_font=args.fallback_font,
            settings_path=args.settings.expanduser().resolve() if args.settings else None,
        )
        data = json.loads(json_path.read_text(encoding="utf-8"))
        if args.json_output:
            copy_json_with_external_assets(json_path, args.json_output.expanduser().resolve())
        emit({
            "ok": True,
            "command": "convert",
            "pdf": str(pdf),
            "pptx": str(pptx),
            "json": str(args.json_output.expanduser().resolve()) if args.json_output else str(json_path),
            "layout": summarize_layout(data),
            "components": list_components(data),
            "pptx_summary": summarize_pptx(pptx),
        })
        return 0
    finally:
        if cleanup:
            cleanup.cleanup()


def command_verify(args: argparse.Namespace) -> int:
    pptx = args.pptx.expanduser().resolve()
    if not pptx.is_file():
        raise FileNotFoundError(f"PPTX not found: {pptx}")
    payload: dict[str, Any] = {
        "ok": True,
        "command": "verify",
        "pptx": str(pptx),
        "pptx_summary": summarize_pptx(pptx),
    }
    if args.json:
        data = json.loads(args.json.expanduser().resolve().read_text(encoding="utf-8"))
        payload["json"] = str(args.json.expanduser().resolve())
        payload["layout"] = summarize_layout(data)
    emit(payload)
    return 0


def command_app(args: argparse.Namespace) -> int:
    desktop = Path(__file__).resolve().parents[1] / "apps" / "desktop"
    if not desktop.is_dir():
        raise FileNotFoundError("Desktop app scaffold was not found at apps/desktop")
    if not shutil.which("npm"):
        raise RuntimeError("npm was not found. Install Node.js or run the desktop app from a packaged build.")
    import subprocess

    completed = subprocess.run(["npm", "run", "tauri", "--", "dev"], cwd=desktop, check=False)
    return completed.returncode


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(prog="pdfppt", description=__doc__)
    subparsers = parser.add_subparsers(dest="command", required=True)

    inspect_parser = subparsers.add_parser("inspect", help="Extract parser JSON and summarize PDF layout.")
    inspect_parser.add_argument("pdf", type=Path)
    inspect_parser.add_argument("--json-output", type=Path, help="Keep parser JSON and image assets at this path.")
    build_parser_common(inspect_parser)
    inspect_parser.set_defaults(func=command_inspect)

    convert_parser = subparsers.add_parser("convert", help="Convert PDF to editable PPTX.")
    convert_parser.add_argument("pdf", type=Path)
    convert_parser.add_argument("pptx", type=Path)
    convert_parser.add_argument("--json-output", type=Path, help="Keep parser JSON and image assets at this path.")
    convert_parser.add_argument("--page-width", type=float, help="Source page width in PDF points.")
    convert_parser.add_argument("--page-height", type=float, help="Source page height in PDF points.")
    convert_parser.add_argument("--fallback-font", default="Pretendard")
    convert_parser.add_argument("--settings", type=Path, help="Optional component settings JSON.")
    build_parser_common(convert_parser)
    convert_parser.set_defaults(func=command_convert)

    verify_parser = subparsers.add_parser("verify", help="Summarize generated PPTX editability.")
    verify_parser.add_argument("pptx", type=Path)
    verify_parser.add_argument("--json", type=Path, help="Optional parser JSON to include in the report.")
    verify_parser.set_defaults(func=command_verify)

    app_parser = subparsers.add_parser("app", help="Launch the local desktop app in development mode.")
    app_parser.set_defaults(func=command_app)

    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        return int(args.func(args) or 0)
    except Exception as exc:
        emit({"ok": False, "command": getattr(args, "command", None), "error": str(exc)})
        return 1
