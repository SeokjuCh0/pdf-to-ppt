#!/usr/bin/env python3
"""Render macOS Vision OCR results into editable PPTX text boxes."""

from __future__ import annotations

import argparse
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


EMU_PER_PT = 12700
DEFAULT_DPI = 144
DEFAULT_FONT = "Pretendard"
DEFAULT_TEXT_COLOR = RGBColor(64, 64, 64)


@dataclass(frozen=True)
class OcrText:
    text: str
    x: float
    y: float
    width: float
    height: float
    confidence: float

    @property
    def font_size(self) -> float:
        return max(4.0, min(24.0, self.height * 0.86))


@dataclass(frozen=True)
class OcrPage:
    image_path: Path
    width_pt: float
    height_pt: float
    texts: list[OcrText]


def pt(value: float) -> Emu:
    return Emu(int(round(value * EMU_PER_PT)))


def parse_pages_range(pages: str | None) -> tuple[int, int]:
    if not pages:
        return 1, 1
    first_part = pages.split(",", 1)[0].strip()
    if "-" in first_part:
        start, end = first_part.split("-", 1)
        return max(1, int(start)), max(1, int(end))
    page = max(1, int(first_part))
    return page, page


def render_pdf_pages(pdf_path: Path, output_dir: Path, pages: str | None, dpi: int, pdftoppm: str) -> list[Path]:
    first, last = parse_pages_range(pages)
    prefix = output_dir / "page"
    command = [
        pdftoppm,
        "-f",
        str(first),
        "-l",
        str(last),
        "-r",
        str(dpi),
        "-png",
        str(pdf_path),
        str(prefix),
    ]
    completed = subprocess.run(command, check=False, capture_output=True, text=True)
    if completed.returncode != 0:
        output = "\n".join(part for part in [completed.stdout, completed.stderr] if part).strip()
        raise RuntimeError(f"pdftoppm failed with exit code {completed.returncode}.\n{output}".strip())
    return sorted(output_dir.glob("page-*.png"))


def vision_frameworks() -> tuple[Any, Any]:
    try:
        import Foundation  # type: ignore
        import Vision  # type: ignore
    except Exception as exc:
        raise RuntimeError("macOS Vision OCR requires PyObjC Foundation and Vision modules.") from exc
    return Foundation, Vision


def recognize_image_text(image_path: Path, width_pt: float, height_pt: float, languages: list[str]) -> list[OcrText]:
    Foundation, Vision = vision_frameworks()
    url = Foundation.NSURL.fileURLWithPath_(str(image_path))
    handler = Vision.VNImageRequestHandler.alloc().initWithURL_options_(url, {})
    request = Vision.VNRecognizeTextRequest.alloc().init()
    request.setRecognitionLevel_(Vision.VNRequestTextRecognitionLevelAccurate)
    request.setRecognitionLanguages_(languages)
    request.setUsesLanguageCorrection_(False)
    ok, error = handler.performRequests_error_([request], None)
    if not ok:
        raise RuntimeError(f"Vision OCR failed: {error}")

    texts: list[OcrText] = []
    for observation in request.results() or []:
        candidates = observation.topCandidates_(1)
        if not candidates:
            continue
        candidate = candidates[0]
        text = str(candidate.string()).strip()
        if not text:
            continue
        box = observation.boundingBox()
        x = float(box.origin.x) * width_pt
        y = (1.0 - float(box.origin.y) - float(box.size.height)) * height_pt
        texts.append(OcrText(
            text=text,
            x=x,
            y=y,
            width=float(box.size.width) * width_pt,
            height=float(box.size.height) * height_pt,
            confidence=float(candidate.confidence()),
        ))
    return texts


def ocr_rendered_pages(page_images: list[Path], dpi: int, languages: list[str]) -> list[OcrPage]:
    pages: list[OcrPage] = []
    for image_path in page_images:
        with Image.open(image_path) as image:
            width_px, height_px = image.size
        width_pt = width_px * 72.0 / dpi
        height_pt = height_px * 72.0 / dpi
        pages.append(OcrPage(
            image_path=image_path,
            width_pt=width_pt,
            height_pt=height_pt,
            texts=recognize_image_text(image_path, width_pt, height_pt, languages),
        ))
    return pages


def set_east_asian_font(run: Any, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = OxmlElement(f"a:{tag}")
            rpr.append(node)
        node.set("typeface", font_name)


def add_ocr_text(slide: Any, text: OcrText, font: str) -> None:
    shape = slide.shapes.add_textbox(pt(text.x), pt(text.y), pt(max(1.0, text.width + 18)), pt(max(1.0, text.height + 3)))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = pt(0)
    frame.margin_right = pt(0)
    frame.margin_top = pt(0)
    frame.margin_bottom = pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text.text
    set_east_asian_font(run, font)
    run.font.size = Pt(text.font_size)
    run.font.color.rgb = DEFAULT_TEXT_COLOR


def render_ocr_pages_to_pptx(pages: list[OcrPage], pptx_path: Path, font: str = DEFAULT_FONT) -> Path:
    if not pages:
        raise ValueError("No OCR pages were produced.")
    prs = Presentation()
    prs.slide_width = pt(pages[0].width_pt)
    prs.slide_height = pt(pages[0].height_pt)
    blank_layout = prs.slide_layouts[6]
    for page in pages:
        slide = prs.slides.add_slide(blank_layout)
        for text in page.texts:
            add_ocr_text(slide, text, font)
    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    return pptx_path


def convert_macos_vision_ocr_to_pptx(
    pdf_path: Path,
    pptx_path: Path,
    *,
    pages: str | None = None,
    dpi: int = DEFAULT_DPI,
    pdftoppm: str = "pdftoppm",
    font: str = DEFAULT_FONT,
    languages: list[str] | None = None,
) -> Path:
    pdf_path = pdf_path.expanduser().resolve()
    pptx_path = pptx_path.expanduser().resolve()
    if not pdf_path.is_file():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    with tempfile.TemporaryDirectory(prefix="pdfppt-vision-ocr-") as tmp:
        images = render_pdf_pages(pdf_path, Path(tmp), pages, dpi, pdftoppm)
        ocr_pages = ocr_rendered_pages(images, dpi, languages or ["ko-KR", "en-US"])
        return render_ocr_pages_to_pptx(ocr_pages, pptx_path, font=font)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pdf", type=Path)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--pages", help='One Poppler page range, for example "1" or "3-5".')
    parser.add_argument("--dpi", type=int, default=DEFAULT_DPI)
    parser.add_argument("--pdftoppm", default="pdftoppm")
    parser.add_argument("--font", default=DEFAULT_FONT)
    parser.add_argument("--languages", default="ko-KR,en-US")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        print(convert_macos_vision_ocr_to_pptx(
            args.pdf,
            args.pptx,
            pages=args.pages,
            dpi=args.dpi,
            pdftoppm=args.pdftoppm,
            font=args.font,
            languages=[value.strip() for value in args.languages.split(",") if value.strip()],
        ))
        return 0
    except Exception as exc:
        print(f"macos_vision_ocr_to_pptx: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
