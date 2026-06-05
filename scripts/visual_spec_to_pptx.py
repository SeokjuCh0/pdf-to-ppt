#!/usr/bin/env python3
"""Render a vision-produced component spec into an editable PPTX."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


EMU_PER_PT = 12700
DEFAULT_WIDTH_IN = 13.333
DEFAULT_HEIGHT_IN = 7.5
DASH_STYLES = {
    "dash": MSO_LINE_DASH_STYLE.DASH,
    "sysdash": MSO_LINE_DASH_STYLE.DASH,
    "sysdot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "dot": MSO_LINE_DASH_STYLE.ROUND_DOT,
    "longdash": MSO_LINE_DASH_STYLE.LONG_DASH,
    "dashdot": MSO_LINE_DASH_STYLE.DASH_DOT,
}


def color(value: Any, default: str | None = None) -> RGBColor | None:
    if value is None:
        value = default
    if value is None:
        return None
    if isinstance(value, str):
        text = value.strip()
        if text.lower() in {"none", "transparent"}:
            return None
        if re.fullmatch(r"#[0-9a-fA-F]{6}", text):
            return RGBColor(int(text[1:3], 16), int(text[3:5], 16), int(text[5:7], 16))
    if isinstance(value, list) and len(value) >= 3:
        channels = [float(v) for v in value[:3]]
        if max(channels) <= 1:
            channels = [v * 255 for v in channels]
        return RGBColor(*(max(0, min(255, round(v))) for v in channels))
    return None


def emu(points: float) -> int:
    return int(round(points * EMU_PER_PT))


def set_fill(shape: Any, fill_value: Any) -> None:
    fill_color = color(fill_value)
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color


def set_line(shape: Any, stroke: Any, width: Any = None, dash: Any = None) -> None:
    stroke_color = color(stroke)
    if stroke_color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = stroke_color
    try:
        shape.line.width = Pt(float(width or 1))
    except (TypeError, ValueError):
        shape.line.width = Pt(1)
    dash_key = str(dash or "").lower().replace("_", "")
    if dash_key in DASH_STYLES:
        shape.line.dash_style = DASH_STYLES[dash_key]


def set_run_font(run: Any, font_family: str, size: Any, text_color: Any, bold: bool = False) -> None:
    run.font.name = font_family
    try:
        run.font.size = Pt(float(size))
    except (TypeError, ValueError):
        run.font.size = Pt(12)
    run.font.bold = bool(bold)
    rgb = color(text_color, "#000000")
    if rgb:
        run.font.color.rgb = rgb
    rpr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = OxmlElement(f"a:{tag}")
            rpr.append(node)
        node.set("typeface", font_family)


def canvas_scale(spec: dict[str, Any]) -> tuple[float, float, float, float]:
    canvas = spec.get("canvas") if isinstance(spec.get("canvas"), dict) else {}
    canvas_width = float(canvas.get("width") or 1920)
    canvas_height = float(canvas.get("height") or 1080)
    width_in = float(canvas.get("width_in") or spec.get("width_in") or DEFAULT_WIDTH_IN)
    height_in = float(canvas.get("height_in") or spec.get("height_in") or DEFAULT_HEIGHT_IN)
    slide_width_pt = width_in * 72
    slide_height_pt = height_in * 72
    return slide_width_pt, slide_height_pt, slide_width_pt / canvas_width, slide_height_pt / canvas_height


def bounds(component: dict[str, Any], sx: float, sy: float) -> tuple[int, int, int, int]:
    x = float(component.get("x", 0)) * sx
    y = float(component.get("y", 0)) * sy
    width = max(1.0, float(component.get("w", component.get("width", 1))) * sx)
    height = max(1.0, float(component.get("h", component.get("height", 1))) * sy)
    return emu(x), emu(y), emu(width), emu(height)


def font_size_points(component: dict[str, Any], sy: float) -> float:
    try:
        if component.get("font_size_px") is not None:
            return float(component["font_size_px"]) * sy
        value = component.get("font_size") or 12
        if str(component.get("font_size_unit") or "").lower() == "px":
            return float(value) * sy
        return float(value)
    except (TypeError, ValueError):
        return 12


def add_text(slide: Any, component: dict[str, Any], sx: float, sy: float) -> None:
    left, top, width, height = bounds(component, sx, sy)
    shape = slide.shapes.add_textbox(left, top, width + emu(float(component.get("extra_right_margin_pt", 60))), height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = bool(component.get("word_wrap", False))
    frame.margin_left = Pt(float(component.get("margin_left_pt", 0)))
    frame.margin_right = Pt(float(component.get("margin_right_pt", 0)))
    frame.margin_top = Pt(float(component.get("margin_top_pt", 0)))
    frame.margin_bottom = Pt(float(component.get("margin_bottom_pt", 0)))
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = str(component.get("text") or "")
    set_run_font(
        run,
        str(component.get("font") or component.get("font_family") or "Arial"),
        font_size_points(component, sy),
        component.get("color") or component.get("text_color") or "#000000",
        bool(component.get("bold", False)),
    )


def add_rect(slide: Any, component: dict[str, Any], sx: float, sy: float) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, *bounds(component, sx, sy))
    set_fill(shape, component.get("fill"))
    set_line(shape, component.get("stroke"), component.get("stroke_width"), component.get("dash"))


def add_line(slide: Any, component: dict[str, Any], sx: float, sy: float) -> None:
    x1 = emu(float(component.get("x1", component.get("x", 0))) * sx)
    y1 = emu(float(component.get("y1", component.get("y", 0))) * sy)
    x2 = emu(float(component.get("x2", component.get("x", 0))) * sx)
    y2 = emu(float(component.get("y2", component.get("y", 0))) * sy)
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    set_line(shape, component.get("stroke"), component.get("stroke_width"), component.get("dash"))


def add_image(slide: Any, component: dict[str, Any], sx: float, sy: float, base_dir: Path) -> None:
    source = component.get("source")
    if not isinstance(source, str):
        return
    image_path = (base_dir / source).resolve()
    if image_path.is_file():
        slide.shapes.add_picture(str(image_path), *bounds(component, sx, sy))


def render_visual_spec(spec_path: Path, pptx_path: Path) -> Path:
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    slide_width_pt, slide_height_pt, sx, sy = canvas_scale(spec)
    prs = Presentation()
    prs.slide_width = Inches(slide_width_pt / 72)
    prs.slide_height = Inches(slide_height_pt / 72)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    base_dir = spec_path.parent

    for component in spec.get("components", []):
        if not isinstance(component, dict) or component.get("include") is False:
            continue
        kind = str(component.get("type") or "").lower().replace("_", "-")
        if kind in {"text", "textbox", "label"}:
            add_text(slide, component, sx, sy)
        elif kind in {"rect", "rectangle", "box"}:
            add_rect(slide, component, sx, sy)
        elif kind in {"line", "connector"}:
            add_line(slide, component, sx, sy)
        elif kind in {"image", "picture"}:
            add_image(slide, component, sx, sy, base_dir)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    return pptx_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("spec", type=Path)
    parser.add_argument("pptx", type=Path)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        print(render_visual_spec(args.spec.expanduser().resolve(), args.pptx.expanduser().resolve()))
        return 0
    except Exception as exc:
        print(f"visual_spec_to_pptx: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
