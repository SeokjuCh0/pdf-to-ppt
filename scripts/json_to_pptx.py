#!/usr/bin/env python3
"""Convert OpenDataLoader-style PDF JSON geometry into an editable PPTX."""

from __future__ import annotations

import argparse
import json
import math
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Pt


EMU_PER_PT = 12700
DEFAULT_PAGE_WIDTH_PT = 612.0
DEFAULT_PAGE_HEIGHT_PT = 792.0
TEXT_TYPES = {"heading", "paragraph", "text chunk", "text block", "list item", "toc item", "formula"}
IMAGE_TYPES = {"image"}
TABLE_TYPES = {"table"}
CJK_FALLBACK = "Pretendard"


@dataclass(frozen=True)
class Rect:
    x: float
    y: float
    width: float
    height: float


def pt(value: float) -> Emu:
    return Emu(int(round(value * EMU_PER_PT)))


def strip_subset_prefix(font_name: str | None) -> str:
    if not font_name:
        return CJK_FALLBACK
    return re.sub(r"^[A-Z]{6}\+", "", str(font_name)).strip() or CJK_FALLBACK


def contains_cjk(text: str) -> bool:
    return any("\uac00" <= ch <= "\ud7a3" or "\u4e00" <= ch <= "\u9fff" or "\u3040" <= ch <= "\u30ff" for ch in text)


def choose_font(font_name: str | None, text: str, fallback: str = CJK_FALLBACK) -> str:
    cleaned = strip_subset_prefix(font_name)
    if cleaned == CJK_FALLBACK and contains_cjk(text):
        return fallback
    return cleaned


def parse_color(value: Any) -> RGBColor | None:
    if value is None:
        return None
    if isinstance(value, str):
        text = value.strip()
        if text.startswith("#") and len(text) == 7:
            try:
                return RGBColor(int(text[1:3], 16), int(text[3:5], 16), int(text[5:7], 16))
            except ValueError:
                return None
        numbers = re.findall(r"-?\d+(?:\.\d+)?", text)
        if len(numbers) >= 3:
            value = [float(n) for n in numbers[:3]]
        else:
            return None
    if isinstance(value, (list, tuple)) and len(value) >= 3:
        channels = [float(v) for v in value[:3]]
        if max(channels) <= 1.0:
            channels = [v * 255 for v in channels]
        return RGBColor(*(max(0, min(255, int(round(v)))) for v in channels))
    return None


def bbox_to_rect(bbox: list[float], page_height: float, scale: float = 1.0) -> Rect:
    left, bottom, right, top = [float(v) for v in bbox]
    x = left * scale
    y = (page_height - top) * scale
    width = max(1.0, (right - left) * scale)
    height = max(1.0, (top - bottom) * scale)
    return Rect(x, y, width, height)


def walk_objects(items: Iterable[dict[str, Any]]) -> Iterable[dict[str, Any]]:
    for item in items:
        if not isinstance(item, dict):
            continue
        yield item
        for key in ("kids", "list items", "toc items"):
            children = item.get(key)
            if isinstance(children, list):
                yield from walk_objects(children)
        rows = item.get("rows")
        if isinstance(rows, list):
            for row in rows:
                cells = row.get("cells") if isinstance(row, dict) else None
                if isinstance(cells, list):
                    yield from walk_objects(cells)


def text_from_object(obj: dict[str, Any]) -> str:
    content = obj.get("content")
    if isinstance(content, str) and content:
        return content
    parts: list[str] = []
    for child in walk_objects(obj.get("kids", [])):
        child_content = child.get("content")
        if isinstance(child_content, str) and child_content:
            parts.append(child_content)
    return "\n".join(parts)


def first_text_style(obj: dict[str, Any]) -> dict[str, Any]:
    if obj.get("font") or obj.get("font size") or obj.get("text color"):
        return obj
    for child in walk_objects(obj.get("kids", [])):
        if child.get("font") or child.get("font size") or child.get("text color"):
            return child
    return obj


def grouped_by_page(data: dict[str, Any]) -> dict[int, list[dict[str, Any]]]:
    grouped: dict[int, list[dict[str, Any]]] = {}
    for obj in data.get("kids", []):
        if not isinstance(obj, dict):
            continue
        page = int(obj.get("page number") or 1)
        grouped.setdefault(page, []).append(obj)
    if not grouped:
        grouped[1] = []
    return grouped


def iter_renderable_objects(items: Iterable[dict[str, Any]]) -> Iterable[dict[str, Any]]:
    for obj in items:
        if not isinstance(obj, dict):
            continue
        obj_type = str(obj.get("type") or "").lower()
        if obj_type in TABLE_TYPES or obj_type in IMAGE_TYPES or obj_type in TEXT_TYPES:
            yield obj
            continue
        for key in ("kids", "list items", "toc items"):
            children = obj.get(key)
            if isinstance(children, list):
                yield from iter_renderable_objects(children)


def metadata_page_size(data: dict[str, Any], page_number: int) -> tuple[float, float] | None:
    pages = data.get("pages")
    if isinstance(pages, list):
        for page in pages:
            if not isinstance(page, dict):
                continue
            number = int(page.get("page number") or page.get("number") or len(pages))
            if number == page_number:
                width = page.get("page width") or page.get("width")
                height = page.get("page height") or page.get("height")
                if width and height:
                    return float(width), float(height)
    width = data.get("page width") or data.get("width")
    height = data.get("page height") or data.get("height")
    if width and height:
        return float(width), float(height)
    return None


def infer_page_size(data: dict[str, Any], page_number: int) -> tuple[float, float]:
    meta = metadata_page_size(data, page_number)
    if meta:
        return meta
    max_right = 0.0
    max_top = 0.0
    for obj in walk_objects(data.get("kids", [])):
        if int(obj.get("page number") or page_number) != page_number:
            continue
        bbox = obj.get("bounding box")
        if isinstance(bbox, list) and len(bbox) == 4:
            max_right = max(max_right, float(bbox[2]))
            max_top = max(max_top, float(bbox[3]))
    return max(max_right, DEFAULT_PAGE_WIDTH_PT), max(max_top, DEFAULT_PAGE_HEIGHT_PT)


def set_run_style(run: Any, style: dict[str, Any], text: str, font_scale: float, fallback_font: str) -> None:
    font_name = choose_font(style.get("font"), text, fallback_font)
    run.font.name = font_name
    size = style.get("font size") or 10
    try:
        run.font.size = Pt(max(4.0, float(size) * font_scale))
    except (TypeError, ValueError):
        run.font.size = Pt(10)
    lower = font_name.lower()
    run.font.bold = any(token in lower for token in ("bold", "semibold", "semi-bold", "medium", "black"))
    run.font.italic = "italic" in lower or "oblique" in lower
    color = parse_color(style.get("text color"))
    if color:
        run.font.color.rgb = color


def add_text_box(slide: Any, obj: dict[str, Any], page_height: float, scale: float, font_scale: float, fallback_font: str) -> None:
    bbox = obj.get("bounding box")
    if not isinstance(bbox, list) or len(bbox) != 4:
        return
    text = text_from_object(obj)
    if not text:
        return
    rect = bbox_to_rect(bbox, page_height, scale)
    shape = slide.shapes.add_textbox(pt(rect.x), pt(rect.y), pt(rect.width), pt(rect.height))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.margin_left = pt(1.5)
    text_frame.margin_right = pt(1.5)
    text_frame.margin_top = pt(0.5)
    text_frame.margin_bottom = pt(0.5)
    paragraph = text_frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, first_text_style(obj), text, font_scale, fallback_font)


def sorted_unique(values: Iterable[float]) -> list[float]:
    rounded = sorted({round(float(v), 3) for v in values})
    return rounded


def table_grid_from_cells(table_obj: dict[str, Any]) -> tuple[list[float], list[float]]:
    x_values: list[float] = []
    y_values: list[float] = []
    for row in table_obj.get("rows", []):
        for cell in row.get("cells", []) if isinstance(row, dict) else []:
            bbox = cell.get("bounding box")
            if isinstance(bbox, list) and len(bbox) == 4:
                left, bottom, right, top = [float(v) for v in bbox]
                x_values.extend([left, right])
                y_values.extend([bottom, top])
    return sorted_unique(x_values), sorted_unique(y_values)


def add_table(slide: Any, table_obj: dict[str, Any], page_height: float, scale: float, font_scale: float, fallback_font: str) -> None:
    bbox = table_obj.get("bounding box")
    if not isinstance(bbox, list) or len(bbox) != 4:
        return
    rows_count = int(table_obj.get("number of rows") or len(table_obj.get("rows", [])) or 1)
    cols_count = int(table_obj.get("number of columns") or 1)
    rect = bbox_to_rect(bbox, page_height, scale)
    shape = slide.shapes.add_table(rows_count, cols_count, pt(rect.x), pt(rect.y), pt(rect.width), pt(rect.height))
    ppt_table = shape.table

    x_grid, y_grid = table_grid_from_cells(table_obj)
    if len(x_grid) >= cols_count + 1:
        for idx in range(cols_count):
            ppt_table.columns[idx].width = pt((x_grid[idx + 1] - x_grid[idx]) * scale)
    if len(y_grid) >= rows_count + 1:
        row_heights = list(reversed([y_grid[idx + 1] - y_grid[idx] for idx in range(len(y_grid) - 1)]))
        for idx, height in enumerate(row_heights[:rows_count]):
            ppt_table.rows[idx].height = pt(height * scale)

    for row in table_obj.get("rows", []):
        if not isinstance(row, dict):
            continue
        for cell in row.get("cells", []):
            r = int(cell.get("row number") or 1) - 1
            c = int(cell.get("column number") or 1) - 1
            if r < 0 or c < 0 or r >= rows_count or c >= cols_count:
                continue
            target = ppt_table.cell(r, c)
            row_span = max(1, int(cell.get("row span") or 1))
            col_span = max(1, int(cell.get("column span") or 1))
            if row_span > 1 or col_span > 1:
                end_r = min(rows_count - 1, r + row_span - 1)
                end_c = min(cols_count - 1, c + col_span - 1)
                try:
                    target = target.merge(ppt_table.cell(end_r, end_c))
                except ValueError:
                    pass
            text = text_from_object(cell)
            target.text = text
            target.margin_left = pt(1)
            target.margin_right = pt(1)
            target.margin_top = pt(0.5)
            target.margin_bottom = pt(0.5)
            if text and target.text_frame.paragraphs:
                paragraph = target.text_frame.paragraphs[0]
                if paragraph.runs:
                    set_run_style(paragraph.runs[0], first_text_style(cell), text, font_scale, fallback_font)


def add_image_or_placeholder(slide: Any, obj: dict[str, Any], page_height: float, scale: float, base_dir: Path | None) -> None:
    bbox = obj.get("bounding box")
    if not isinstance(bbox, list) or len(bbox) != 4:
        return
    rect = bbox_to_rect(bbox, page_height, scale)
    source = obj.get("source")
    image_path = (base_dir / source).resolve() if base_dir and isinstance(source, str) else None
    if image_path and image_path.is_file():
        slide.shapes.add_picture(str(image_path), pt(rect.x), pt(rect.y), pt(rect.width), pt(rect.height))
        return
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pt(rect.x), pt(rect.y), pt(rect.width), pt(rect.height))
    shape.text = obj.get("alt") or "image"


def convert_json_to_pptx(
    json_path: Path,
    pptx_path: Path,
    *,
    page_width: float | None = None,
    page_height: float | None = None,
    fallback_font: str = CJK_FALLBACK,
) -> Path:
    data = json.loads(json_path.read_text(encoding="utf-8"))
    grouped = grouped_by_page(data)
    first_page = min(grouped)
    inferred_width, inferred_height = infer_page_size(data, first_page)
    slide_width = float(page_width or inferred_width)
    slide_height = float(page_height or inferred_height)

    prs = Presentation()
    prs.slide_width = pt(slide_width)
    prs.slide_height = pt(slide_height)
    blank_layout = prs.slide_layouts[6]
    base_dir = json_path.parent

    for page in sorted(grouped):
        native_width, native_height = infer_page_size(data, page)
        if page_width and page_height:
            native_width, native_height = float(page_width), float(page_height)
        scale = min(slide_width / native_width, slide_height / native_height) if native_width and native_height else 1.0
        font_scale = scale
        slide = prs.slides.add_slide(blank_layout)
        for obj in iter_renderable_objects(grouped[page]):
            obj_type = str(obj.get("type") or "").lower()
            if obj_type in TABLE_TYPES:
                add_table(slide, obj, native_height, scale, font_scale, fallback_font)
            elif obj_type in IMAGE_TYPES:
                add_image_or_placeholder(slide, obj, native_height, scale, base_dir)
            elif obj_type in TEXT_TYPES:
                add_text_box(slide, obj, native_height, scale, font_scale, fallback_font)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    return pptx_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("json", type=Path, help="OpenDataLoader JSON path.")
    parser.add_argument("pptx", type=Path, help="Output PPTX path.")
    parser.add_argument("--page-width", type=float, help="Source page width in PDF points.")
    parser.add_argument("--page-height", type=float, help="Source page height in PDF points.")
    parser.add_argument("--fallback-font", default=CJK_FALLBACK, help="Fallback font for missing/CJK fonts.")
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        output = convert_json_to_pptx(
            args.json.expanduser().resolve(),
            args.pptx.expanduser().resolve(),
            page_width=args.page_width,
            page_height=args.page_height,
            fallback_font=args.fallback_font,
        )
        print(output)
        return 0
    except Exception as exc:
        print(f"json_to_pptx: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
