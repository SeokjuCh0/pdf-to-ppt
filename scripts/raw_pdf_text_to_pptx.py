#!/usr/bin/env python3
"""Render Poppler PDF text bbox output into editable PPTX text boxes."""

from __future__ import annotations

import argparse
import subprocess
import sys
import tempfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


EMU_PER_PT = 12700
DEFAULT_FONT = "Pretendard"
DEFAULT_TEXT_COLOR = RGBColor(64, 64, 64)


@dataclass(frozen=True)
class TextLine:
    page_index: int
    page_width: float
    page_height: float
    x_min: float
    y_min: float
    x_max: float
    y_max: float
    text: str

    @property
    def font_size(self) -> float:
        return max(4.0, min(28.0, (self.y_max - self.y_min) * 0.88))


def pt(value: float) -> Emu:
    return Emu(int(round(value * EMU_PER_PT)))


def tag_name(element: ET.Element) -> str:
    return element.tag.rsplit("}", 1)[-1]


def float_attr(element: ET.Element, name: str, default: float = 0.0) -> float:
    try:
        return float(element.attrib.get(name, default))
    except (TypeError, ValueError):
        return default


def iter_elements(root: ET.Element, name: str) -> Iterable[ET.Element]:
    for element in root.iter():
        if tag_name(element) == name:
            yield element


def line_text(line: ET.Element) -> str:
    words: list[str] = []
    for word in line:
        if tag_name(word) != "word":
            continue
        value = "".join(word.itertext()).strip()
        if value:
            words.append(value)
    return " ".join(words)


def parse_bbox_html(html_path: Path) -> list[TextLine]:
    root = ET.parse(html_path).getroot()
    lines: list[TextLine] = []
    for page_index, page in enumerate(iter_elements(root, "page"), start=1):
        page_width = float_attr(page, "width", 612.0)
        page_height = float_attr(page, "height", 792.0)
        for line in iter_elements(page, "line"):
            text = line_text(line)
            if not text:
                continue
            lines.append(TextLine(
                page_index=page_index,
                page_width=page_width,
                page_height=page_height,
                x_min=float_attr(line, "xMin"),
                y_min=float_attr(line, "yMin"),
                x_max=float_attr(line, "xMax"),
                y_max=float_attr(line, "yMax"),
                text=text,
            ))
    return lines


def set_east_asian_font(run: Any, font_name: str) -> None:
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        node = rpr.find(qn(f"a:{tag}"))
        if node is None:
            node = OxmlElement(f"a:{tag}")
            rpr.append(node)
        node.set("typeface", font_name)


def add_text_line(slide: Any, line: TextLine, scale: float, font: str) -> None:
    width = max(1.0, (line.x_max - line.x_min) * scale)
    height = max(1.0, (line.y_max - line.y_min) * scale)
    shape = slide.shapes.add_textbox(
        pt(line.x_min * scale),
        pt(line.y_min * scale),
        pt(width + 48),
        pt(height + 4),
    )
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = pt(0)
    frame.margin_right = pt(0)
    frame.margin_top = pt(0)
    frame.margin_bottom = pt(0)
    paragraph = frame.paragraphs[0]
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = line.text
    set_east_asian_font(run, font)
    run.font.size = Pt(line.font_size * scale)
    run.font.color.rgb = DEFAULT_TEXT_COLOR


def render_text_lines_to_pptx(lines: list[TextLine], pptx_path: Path, font: str = DEFAULT_FONT) -> Path:
    if not lines:
        raise ValueError("No PDF text lines were extracted.")

    pages: dict[int, list[TextLine]] = {}
    for line in lines:
        pages.setdefault(line.page_index, []).append(line)

    first_page = pages[min(pages)][0]
    prs = Presentation()
    prs.slide_width = pt(first_page.page_width)
    prs.slide_height = pt(first_page.page_height)
    blank_layout = prs.slide_layouts[6]

    for page_index in sorted(pages):
        page_lines = pages[page_index]
        slide = prs.slides.add_slide(blank_layout)
        page_width = page_lines[0].page_width
        page_height = page_lines[0].page_height
        scale = min(first_page.page_width / page_width, first_page.page_height / page_height)
        for line in page_lines:
            add_text_line(slide, line, scale, font)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)
    return pptx_path


def extract_bbox_html(pdf_path: Path, html_path: Path, pages: str | None = None, pdftotext: str = "pdftotext") -> Path:
    command = [pdftotext]
    if pages:
        first, last = parse_pages_range(pages)
        command.extend(["-f", str(first), "-l", str(last)])
    command.extend(["-bbox-layout", str(pdf_path), str(html_path)])
    completed = subprocess.run(command, check=False, capture_output=True, text=True)
    if completed.returncode != 0:
        output = "\n".join(part for part in [completed.stdout, completed.stderr] if part).strip()
        raise RuntimeError(f"pdftotext failed with exit code {completed.returncode}.\n{output}".strip())
    return html_path


def parse_pages_range(pages: str) -> tuple[int, int]:
    first_part = pages.split(",", 1)[0].strip()
    if "-" in first_part:
        start, end = first_part.split("-", 1)
        return max(1, int(start)), max(1, int(end))
    page = max(1, int(first_part))
    return page, page


def convert_raw_pdf_text_to_pptx(
    pdf_path: Path,
    pptx_path: Path,
    *,
    pages: str | None = None,
    pdftotext: str = "pdftotext",
    font: str = DEFAULT_FONT,
) -> Path:
    pdf_path = pdf_path.expanduser().resolve()
    pptx_path = pptx_path.expanduser().resolve()
    if not pdf_path.is_file():
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    with tempfile.TemporaryDirectory(prefix="pdfppt-bbox-") as tmp:
        html_path = Path(tmp) / "bbox.html"
        extract_bbox_html(pdf_path, html_path, pages=pages, pdftotext=pdftotext)
        return render_text_lines_to_pptx(parse_bbox_html(html_path), pptx_path, font=font)


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pdf", type=Path)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--pages", help='One Poppler page range, for example "1" or "3-5".')
    parser.add_argument("--pdftotext", default="pdftotext")
    parser.add_argument("--font", default=DEFAULT_FONT)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        print(convert_raw_pdf_text_to_pptx(
            args.pdf,
            args.pptx,
            pages=args.pages,
            pdftotext=args.pdftotext,
            font=args.font,
        ))
        return 0
    except Exception as exc:
        print(f"raw_pdf_text_to_pptx: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
