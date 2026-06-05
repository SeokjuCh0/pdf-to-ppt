from pathlib import Path
import sys
import tempfile
import unittest

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from macos_vision_ocr_to_pptx import OcrPage, OcrText, render_ocr_pages_to_pptx


class MacosVisionOcrToPptxTests(unittest.TestCase):
    def test_ocr_results_render_as_editable_text(self):
        page = OcrPage(
            image_path=Path("page.png"),
            width_pt=595.32,
            height_pt=841.92,
            texts=[OcrText("42.8", 120, 240, 28, 10, 0.98)],
        )
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "ocr.pptx"
            render_ocr_pages_to_pptx([page], output)

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            texts = [shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("42.8", texts)


if __name__ == "__main__":
    unittest.main()
