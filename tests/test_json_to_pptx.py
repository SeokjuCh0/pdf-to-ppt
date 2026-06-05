from pathlib import Path
import json
import sys
import tempfile
import unittest

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from json_to_pptx import bbox_to_rect, convert_json_to_pptx, strip_subset_prefix
from pdf_to_pptx import copy_json_with_external_assets


class JsonToPptxTests(unittest.TestCase):
    def test_bbox_to_rect_flips_pdf_origin(self):
        rect = bbox_to_rect([20, 90, 380, 210], page_height=300)
        self.assertEqual(rect.x, 20)
        self.assertEqual(rect.y, 90)
        self.assertEqual(rect.width, 360)
        self.assertEqual(rect.height, 120)

    def test_strip_subset_prefix(self):
        self.assertEqual(strip_subset_prefix("ABCDEF+Pretendard-Regular"), "Pretendard-Regular")

    def test_fixture_conversion_creates_editable_shapes(self):
        fixture = ROOT / "tests" / "fixtures" / "opendataloader_sample.json"
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "fixture.pptx"
            convert_json_to_pptx(fixture, output)
            self.assertTrue(output.exists())

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            self.assertEqual(prs.slide_width, 400 * 12700)
            self.assertEqual(prs.slide_height, 300 * 12700)
            texts = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            table_texts = []
            for shape in prs.slides[0].shapes:
                if not shape.has_table:
                    continue
                for row in shape.table.rows:
                    for cell in row.cells:
                        table_texts.append(cell.text)
            texts = "\n".join([texts, *table_texts])
            self.assertIn("2025 한국 부자 보고서", texts)
            self.assertIn("항목", texts)
            self.assertIn("1.0%", texts)

    def test_component_settings_hide_and_override_text(self):
        fixture = ROOT / "tests" / "fixtures" / "opendataloader_sample.json"
        with tempfile.TemporaryDirectory() as tmp:
            settings = Path(tmp) / "settings.json"
            settings.write_text(
                json.dumps({
                    "components": {
                        "1": {"include": False},
                        "2": {"content": "사용자 수정 본문", "font_size": 12},
                    }
                }),
                encoding="utf-8",
            )
            output = Path(tmp) / "fixture.pptx"
            convert_json_to_pptx(fixture, output, settings_path=settings)

            prs = Presentation(output)
            texts = "\n".join(shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text"))
            self.assertNotIn("2025 한국 부자 보고서", texts)
            self.assertIn("사용자 수정 본문", texts)

    def test_nested_list_items_are_rendered_as_text_boxes(self):
        data = {
            "file name": "list.pdf",
            "number of pages": 1,
            "page width": 400,
            "page height": 300,
            "kids": [
                {
                    "type": "list",
                    "page number": 1,
                    "bounding box": [20, 120, 380, 260],
                    "list items": [
                        {
                            "type": "list item",
                            "page number": 1,
                            "bounding box": [20, 200, 380, 240],
                            "font": "ABCDEE+Pretendard-Regular",
                            "font size": 10,
                            "content": "첫 번째 본문 문단",
                        },
                        {
                            "type": "list item",
                            "page number": 1,
                            "bounding box": [20, 130, 380, 170],
                            "font": "ABCDEE+Pretendard-Regular",
                            "font size": 10,
                            "content": "두 번째 본문 문단",
                        },
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            fixture = Path(tmp) / "list.json"
            fixture.write_text(json.dumps(data), encoding="utf-8")
            output = Path(tmp) / "list.pptx"
            convert_json_to_pptx(fixture, output)

            prs = Presentation(output)
            texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
            self.assertIn("첫 번째 본문 문단", texts)
            self.assertIn("두 번째 본문 문단", texts)

    def test_component_settings_apply_to_synthetic_ids(self):
        data = {
            "file name": "list.pdf",
            "number of pages": 1,
            "page width": 400,
            "page height": 300,
            "kids": [
                {
                    "type": "list",
                    "page number": 1,
                    "bounding box": [20, 120, 380, 260],
                    "list items": [
                        {
                            "type": "list item",
                            "page number": 1,
                            "bounding box": [20, 200, 380, 240],
                            "font size": 10,
                            "content": "원본 문단",
                        }
                    ],
                }
            ],
        }
        with tempfile.TemporaryDirectory() as tmp:
            fixture = Path(tmp) / "list.json"
            fixture.write_text(json.dumps(data), encoding="utf-8")
            settings = Path(tmp) / "settings.json"
            settings.write_text(
                json.dumps({"components": {"1:1": {"content": "합성 ID 수정 문단"}}}),
                encoding="utf-8",
            )
            output = Path(tmp) / "list.pptx"
            convert_json_to_pptx(fixture, output, settings_path=settings)

            prs = Presentation(output)
            texts = [shape.text for shape in prs.slides[0].shapes if hasattr(shape, "text")]
            self.assertIn("합성 ID 수정 문단", texts)
            self.assertNotIn("원본 문단", texts)

    def test_json_output_keeps_external_image_assets_reusable(self):
        with tempfile.TemporaryDirectory() as tmp:
            source_dir = Path(tmp) / "source"
            target_dir = Path(tmp) / "target"
            image = source_dir / "fixture_images" / "imageFile1.png"
            image.parent.mkdir(parents=True)
            image.write_bytes(b"fake image")
            json_path = source_dir / "fixture.json"
            json_path.write_text(
                json.dumps({"kids": [{"type": "image", "source": "fixture_images/imageFile1.png"}]}),
                encoding="utf-8",
            )

            target = target_dir / "fixture.json"
            copy_json_with_external_assets(json_path, target)

            self.assertTrue(target.exists())
            self.assertEqual((target_dir / "fixture_images" / "imageFile1.png").read_bytes(), b"fake image")


if __name__ == "__main__":
    unittest.main()
