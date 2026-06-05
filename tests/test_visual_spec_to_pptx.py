import json
from pathlib import Path
import sys
import tempfile
import unittest

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from visual_spec_to_pptx import render_visual_spec


class VisualSpecToPptxTests(unittest.TestCase):
    def test_visual_spec_renders_native_shapes(self):
        fixture = ROOT / "tests" / "fixtures" / "visual_spec_sample.json"
        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "visual.pptx"
            render_visual_spec(fixture, output)

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            slide = prs.slides[0]
            self.assertGreaterEqual(len(slide.shapes), 3)
            texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("Editable visual reconstruction", texts)

    def test_pixel_font_size_is_normalized_to_points(self):
        spec = {
            "canvas": {"width": 1600, "height": 900, "width_in": 13.333, "height_in": 7.5},
            "components": [{
                "type": "text",
                "x": 0,
                "y": 0,
                "w": 300,
                "h": 80,
                "text": "Scaled font",
                "font_size_px": 30,
            }],
        }
        with tempfile.TemporaryDirectory() as tmp:
            spec_path = Path(tmp) / "spec.json"
            output = Path(tmp) / "visual.pptx"
            spec_path.write_text(json.dumps(spec), encoding="utf-8")

            render_visual_spec(spec_path, output)

            prs = Presentation(output)
            run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
            self.assertAlmostEqual(run.font.size.pt, 18.0, places=1)

    def test_chart_component_renders_native_chart(self):
        spec = {
            "canvas": {"width": 1600, "height": 900, "width_in": 13.333, "height_in": 7.5},
            "components": [{
                "type": "chart",
                "chart_type": "stacked-column-100",
                "x": 100,
                "y": 120,
                "w": 900,
                "h": 360,
                "categories": ["2023", "2024"],
                "series": [
                    {"name": "금융", "values": [17.0, 12.0], "color": "#FAC858"},
                    {"name": "부동산", "values": [29.8, 27.0], "color": "#4FA8A5"},
                ],
                "data_labels": True,
            }],
        }
        with tempfile.TemporaryDirectory() as tmp:
            spec_path = Path(tmp) / "spec.json"
            output = Path(tmp) / "visual.pptx"
            spec_path.write_text(json.dumps(spec), encoding="utf-8")

            render_visual_spec(spec_path, output)

            prs = Presentation(output)
            chart_shapes = [shape for shape in prs.slides[0].shapes if getattr(shape, "has_chart", False)]
            self.assertEqual(len(chart_shapes), 1)
            self.assertEqual(len(chart_shapes[0].chart.series), 2)


if __name__ == "__main__":
    unittest.main()
