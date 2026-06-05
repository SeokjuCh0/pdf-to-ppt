from pathlib import Path
import sys
import tempfile
import unittest

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from raw_pdf_text_to_pptx import parse_bbox_html, render_text_lines_to_pptx


class RawPdfTextToPptxTests(unittest.TestCase):
    def test_poppler_bbox_lines_render_as_editable_text(self):
        fixture = ROOT / "tests" / "fixtures" / "poppler_bbox_sample.html"
        lines = parse_bbox_html(fixture)

        self.assertEqual(len(lines), 1)
        self.assertEqual(lines[0].text, "Graph label 42.8")

        with tempfile.TemporaryDirectory() as tmp:
            output = Path(tmp) / "raw-text.pptx"
            render_text_lines_to_pptx(lines, output)

            prs = Presentation(output)
            self.assertEqual(len(prs.slides), 1)
            texts = [shape.text for shape in prs.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("Graph label 42.8", texts)


if __name__ == "__main__":
    unittest.main()
